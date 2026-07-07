"""업로드된 문서에서 검색용 텍스트를 추출한다.

지원 형식: PDF, DOCX, PPTX, XLSX, HWP(미리보기 텍스트), TXT/MD/CSV 등 텍스트 파일.
추출 실패는 검색 불가일 뿐 업로드 실패는 아니므로 예외를 삼키고 빈 문자열을 돌려준다.
"""

import io
import logging
from pathlib import Path

from ..config import MAX_TEXT_LEN

logger = logging.getLogger(__name__)

TEXT_SUFFIXES = {".txt", ".md", ".csv", ".log", ".json", ".xml", ".yaml", ".yml"}


def extract_text(data: bytes, filename: str) -> str:
    suffix = Path(filename).suffix.lower()
    try:
        if suffix == ".pdf":
            text = _from_pdf(data)
        elif suffix == ".docx":
            text = _from_docx(data)
        elif suffix == ".pptx":
            text = _from_pptx(data)
        elif suffix == ".xlsx":
            text = _from_xlsx(data)
        elif suffix == ".hwp":
            text = _from_hwp(data)
        elif suffix in TEXT_SUFFIXES:
            text = data.decode("utf-8", errors="replace")
        else:
            text = ""
    except Exception:
        logger.exception("텍스트 추출 실패: %s", filename)
        text = ""
    return text[:MAX_TEXT_LEN]


def _from_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _from_docx(data: bytes) -> str:
    import docx

    document = docx.Document(io.BytesIO(data))
    parts = [p.text for p in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            parts.append("\t".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def _from_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    return "\n".join(parts)


def _from_xlsx(data: bytes) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"[{ws.title}]")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append("\t".join(cells))
            if sum(len(p) for p in parts) > MAX_TEXT_LEN:
                return "\n".join(parts)
    return "\n".join(parts)


def _from_hwp(data: bytes) -> str:
    """HWP 5.x 본문(BodyText/Section*)을 직접 파싱해 전문을 추출한다.

    본문 파싱이 실패하거나 비어 있으면(암호화·배포용 문서 등) PrvText(미리보기)로 폴백.
    """
    import olefile

    ole = olefile.OleFileIO(io.BytesIO(data))
    try:
        try:
            text = _hwp_body_text(ole)
        except Exception:
            logger.exception("HWP 본문 파싱 실패 — PrvText로 폴백")
            text = ""
        if text.strip():
            return text
        if ole.exists("PrvText"):
            return ole.openstream("PrvText").read().decode("utf-16-le", errors="replace")
        return ""
    finally:
        ole.close()


# 레코드 헤더: uint32 = tag(10bit) | level(10bit) | size(12bit)
_HWPTAG_PARA_TEXT = 0x10 + 51
# 8 WCHAR(16바이트)를 차지하는 인라인/확장 컨트롤 문자 코드
_HWP_CTRL8 = frozenset({1, 2, 3, 4, 5, 6, 7, 8, 9, 11, 12, 14, 15, 16, 17, 18, 19, 20, 21, 22, 23})


def _hwp_body_text(ole) -> str:
    import zlib

    if not ole.exists("FileHeader"):
        return ""
    header = ole.openstream("FileHeader").read()
    flags = int.from_bytes(header[36:40], "little")
    if flags & 0x2:  # 암호화 문서는 본문을 읽을 수 없다
        return ""
    compressed = bool(flags & 0x1)
    sections = sorted(
        (e for e in ole.listdir() if len(e) == 2 and e[0] == "BodyText" and e[1].startswith("Section")),
        key=lambda e: int(e[1][len("Section"):]),
    )
    parts = []
    for entry in sections:
        raw = ole.openstream("/".join(entry)).read()
        if compressed:
            raw = zlib.decompress(raw, -15)
        parts.append(_hwp_section_text(raw))
    return "\n".join(p for p in parts if p)


def _hwp_section_text(data: bytes) -> str:
    parts = []
    pos, end = 0, len(data)
    while pos + 4 <= end:
        header = int.from_bytes(data[pos:pos + 4], "little")
        pos += 4
        tag = header & 0x3FF
        size = (header >> 20) & 0xFFF
        if size == 0xFFF:  # 확장 크기
            size = int.from_bytes(data[pos:pos + 4], "little")
            pos += 4
        if tag == _HWPTAG_PARA_TEXT:
            parts.append(_hwp_para_text(data[pos:pos + size]))
        pos += size
    return "\n".join(p for p in parts if p)


def _hwp_para_text(payload: bytes) -> str:
    """PARA_TEXT 레코드의 UTF-16LE 문자열에서 컨트롤 문자를 걷어낸다."""
    out: list[str] = []
    end = len(payload) - (len(payload) % 2)
    run_start = 0
    i = 0

    def flush(upto: int) -> None:
        if upto > run_start:
            out.append(payload[run_start:upto].decode("utf-16-le", errors="replace"))

    while i < end:
        code = int.from_bytes(payload[i:i + 2], "little")
        if code >= 32:
            i += 2
            continue
        flush(i)
        if code in _HWP_CTRL8:
            if code == 9:
                out.append("\t")
            i += 16
        else:
            if code in (10, 13):
                out.append("\n")
            i += 2
        run_start = i
    flush(end)
    return "".join(out).strip()
